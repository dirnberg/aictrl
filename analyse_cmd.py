import os
import argparse
from pptx import Presentation
import fitz  # PyMuPDF
import openai
import time
import re
from collections import Counter
import string

def debug(message):
    """Print debug messages if debugging is enabled."""
    if DEBUG:
        print(f"[DEBUG] {message}")

def translate_text(text, notes_text="", comments_text=""):
    """Translate slide or PDF text to English using OpenAI."""
    translation_prompt = f"Translate the following text to English:\n\n{text}\n"
    if notes_text:
        translation_prompt += f"\nSpeaker Notes:\n{notes_text}\n"
    if comments_text:
        translation_prompt += f"\nComments:\n{comments_text}\n"

    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that translates and explains content."},
                {"role": "user", "content": translation_prompt}
            ]
        )
        return response.choices[0].message['content'].strip(), response['usage']['total_tokens']
    except Exception as e:
        print(f"Failed to translate text with OpenAI: {e}")
        return text, 0  # Fallback to original text if translation fails

def generate_explanation(translated_text):
    """Generate an explanation for the translated text using OpenAI."""
    explanation_prompt = f"Explain the following translated content as if presenting it:\n\n{translated_text}\n"
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that explains content."},
                {"role": "user", "content": explanation_prompt}
            ]
        )
        return response.choices[0].message['content'].strip(), response['usage']['total_tokens']
    except Exception as e:
        print(f"Failed to generate explanation with OpenAI: {e}")
        return translated_text, 0  # Fallback to translated text if explanation fails

def generate_summary(key_points_en, author):
    """Generate a Blinkist-style summary of the presentation or document using OpenAI."""
    summary_prompt_en = f"Create a Blinkist-style summary of the following key points from a presentation or document by {author}:\n\n" + "\n\n".join(key_points_en)
    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that summarizes content in a Blinkist style."},
                {"role": "user", "content": summary_prompt_en}
            ]
        )
        return response.choices[0].message['content'].strip(), response['usage']['total_tokens']
    except Exception as e:
        print(f"Failed to generate Blinkist-style summary with OpenAI: {e}")
        return "Summary not available due to an error.", 0

def extract_author_from_content(content):
    """Extract the author's name from the content."""
    author_name = None
    author_patterns = [
        r'(Author|Presented by|Created by|Prepared by):?\s*([A-Za-z\s]+)',
        r'By\s([A-Za-z\s]+)'
    ]
    
    for pattern in author_patterns:
        match = re.search(pattern, content, re.IGNORECASE)
        if match:
            author_name = match.group(2).strip()
            break
    
    return author_name or "Unknown Author"

def generate_tags(content):
    """Generate tags based on the most common words in the content."""
    text = content.lower()
    # Remove punctuation and numbers
    text = text.translate(str.maketrans('', '', string.punctuation + string.digits))
    # Split into words and remove common stop words
    words = text.split()
    
    # Manually define stop words
    stop_words = {"a", "an", "the", "and", "or", "but", "if", "then", "is", "are", "was", "were", "to", "from", "in", "on", "with", "as", "by", "for", "of", "at", "this", "that", "these", "those", "it", "its", "they", "them", "their"}
    
    filtered_words = [word for word in words if word not in stop_words and len(word) > 2]
    
    # Count word frequency
    word_counts = Counter(filtered_words)
    most_common_words = [word for word, count in word_counts.most_common(10)]
    
    return ', '.join(most_common_words)

def generate_title_author_tags(content):
    """Generate title, author, and tags based on the content."""
    prompt = f"Based on the following content, provide a suitable title, author, and keywords:\n\n{content}\n\nProvide the response in the following format:\n\nTitle: <Title>\nAuthor: <Author>\nTags: <Comma-separated tags>"

    try:
        response = openai.ChatCompletion.create(
            model="gpt-4-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that generates titles, author names, and keywords for documents."},
                {"role": "user", "content": prompt}
            ]
        )
        response_content = response.choices[0].message['content'].strip()
        title_match = re.search(r'Title:\s*(.*)', response_content)
        author_match = re.search(r'Author:\s*(.*)', response_content)
        tags_match = re.search(r'Tags:\s*(.*)', response_content)
        
        title = title_match.group(1).strip() if title_match else "Untitled Document"
        author = author_match.group(1).strip() if author_match else "Unknown Author"
        tags = tags_match.group(1).strip() if tags_match else "No Tags Available"
        
        return title, author, tags, response['usage']['total_tokens']
    except Exception as e:
        print(f"Failed to generate title, author, and tags with OpenAI: {e}")
        return "Untitled Document", "Unknown Author", "No Tags Available", 0

def save_markdown(filename, content):
    """Save the given content to a Markdown file."""
    with open(filename, "w") as md_file:
        md_file.write(content)

def process_pptx(file_path):
    """Process a PowerPoint file."""
    ppt = Presentation(file_path)
    slide_contents = []
    key_points_en = []
    total_tokens = 0

    for i, slide in enumerate(ppt.slides):
        debug(f"Processing slide {i+1} of {len(ppt.slides)}.")

        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slide_contents.append(slide_text)
        
        notes_text = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        comments_text = "\n".join([f"{comment.author}: {comment.text}" for comment in slide.notes_slide.comments]) if slide.has_notes_slide and hasattr(slide.notes_slide, 'comments') else ""

        translated_text, tokens = translate_text(slide_text, notes_text, comments_text)
        total_tokens += tokens

        slide_explanation, tokens = generate_explanation(translated_text)
        total_tokens += tokens

        key_points_en.append(slide_explanation)
        print(f"Progress: {(i + 1) / len(ppt.slides) * 100:.2f}%")

    return slide_contents, key_points_en, total_tokens

def process_pdf(file_path):
    """Process a PDF file."""
    pdf = fitz.open(file_path)
    content = []
    key_points_en = []
    total_tokens = 0

    for i, page in enumerate(pdf):
        debug(f"Processing page {i+1} of {len(pdf)}.")

        page_text = page.get_text()
        content.append(page_text)
        
        translated_text, tokens = translate_text(page_text)
        total_tokens += tokens

        page_explanation, tokens = generate_explanation(translated_text)
        total_tokens += tokens

        key_points_en.append(page_explanation)
        print(f"Progress: {(i + 1) / len(pdf) * 100:.2f}%")

    return content, key_points_en, total_tokens

if __name__ == "__main__":
    DEBUG = True
    start_time = time.time()
    openai.api_key = os.getenv("OPENAI_API_KEY")

    parser = argparse.ArgumentParser(description='Convert a PowerPoint or PDF file to a Markdown report with English explanations, professional knowledge management summary, and source listing.')
    parser.add_argument('file', type=str, help='The PowerPoint or PDF file to analyze (.pptx or .pdf)')
    args = parser.parse_args()

    file_extension = os.path.splitext(args.file)[1].lower()

    if file_extension == ".pptx":
        slide_contents, key_points_en, total_tokens = process_pptx(args.file)
    elif file_extension == ".pdf":
        slide_contents, key_points_en, total_tokens = process_pdf(args.file)
    else:
        raise ValueError("Unsupported file type. Please provide a .pptx or .pdf file.")

    content = ' '.join(slide_contents)
    title, author, tags, prompt_tokens = generate_title_author_tags(content)
    total_tokens += prompt_tokens

    summary_en, tokens = generate_summary(key_points_en, author)
    total_tokens += tokens

    markdown_output_en = (
        f"# {title}\n\n"
        f"*Author: {author}*\n\n"
        f"*Date: {time.strftime('%Y-%m-%d')}*\n\n"
        f"*Tags: {tags}*\n\n"
        f"## Blinkist-Style Summary\n\n"
        f"{summary_en}\n\n"
        f"## Document Properties\n\n"
        f"- **Title**: {title}\n"
        f"- **Author**: {author}\n"
        f"- **Creation Date**: {time.strftime('%Y-%m-%d')}\n"
        f"- **Keywords**: {tags}\n\n"
        f"## Slide-by-Slide (or Page-by-Page) Analysis\n"
    )

    for i, explanation in enumerate(key_points_en):
        markdown_output_en += (
            f"\n### Page {i+1}\n\n"
            f"**Explanation:**\n\n{explanation}\n"
        )

    save_markdown(os.path.splitext(args.file)[0] + ".md", markdown_output_en)

    cost_per_token = 0.03 / 1000  # Example cost, update according to your pricing
    euro_conversion_rate = 0.85
    total_cost_eur = total_tokens * cost_per_token * euro_conversion_rate
    elapsed_time = time.time() - start_time

    print(f"\n---\nTotal Time Elapsed: {elapsed_time:.2f} seconds")
    print(f"Total Tokens Used: {total_tokens}")
    print(f"Total Cost (EUR): {total_cost_eur:.4f}")
